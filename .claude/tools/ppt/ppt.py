#!/usr/bin/env python3
"""
Odoo 工時表同步與 PPT 自動生成腳本
自動動態計算「執行當天前兩週的週一至週五」，建立目錄並生成左右雙表格 PPT。
"""

import sys
from datetime import datetime, timedelta
from pathlib import Path
from collections import defaultdict

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from _odoo_utils import load_config, remove_images_only, create_odoo_session


def get_previous_two_weeks_range():
    today = datetime.now().date()
    this_week_monday = today - timedelta(days=today.weekday())
    two_weeks_ago_monday = this_week_monday - timedelta(weeks=2)
    last_week_friday = two_weeks_ago_monday + timedelta(days=11)
    return two_weeks_ago_monday.strftime("%Y-%m-%d"), last_week_friday.strftime("%Y-%m-%d")


_A01_TASK_MAP = [
    ("B01", "會議"),
    ("B02", "教育訓練"),
    ("B03", "休假"),
]


def _resolve_pending_key(project_name, task_name):
    """將特定專案的任務細分為子分類 key；其餘直接用專案名稱。"""
    if project_name == "A01-艾創點內部/組織管理":
        for prefix, label in _A01_TASK_MAP:
            if task_name.startswith(prefix):
                return label
        return project_name
    if project_name == "A03-案件eService處理":
        return task_name or project_name
    return project_name


def fetch_pending_tasks(odoo_cfg, service_cfg):
    """
    從兩個來源抓未完成任務，回傳 {project_name: count}。
    service_cfg 為 None 時略過第二來源。
    """
    stats = defaultdict(int)

    # 來源 1：project.task
    session = create_odoo_session(odoo_cfg["url"], odoo_cfg["db"], odoo_cfg["username"], odoo_cfg["password"])
    resp = session.post(f"{odoo_cfg['url']}/web/dataset/call_kw", json={
        "jsonrpc": "2.0",
        "method": "call",
        "params": {
            "model": "project.task",
            "method": "search_read",
            "args": [],
            "kwargs": {
                "domain": [["user_id", "=", odoo_cfg["user_id"]]],
                "fields": ["project_id", "name"],
                "limit": 200,
            },
        },
    }).json()
    if "result" in resp:
        for task in resp["result"]:
            project_name = task["project_id"][1] if task.get("project_id") else "未知專案"
            task_name = task.get("name", "")
            key = _resolve_pending_key(project_name, task_name)
            stats[key] += 1
    else:
        print(f"[WARN] project.task 查詢失敗: {resp.get('error')}")

    # 來源 2：service.question.feedback（選用）
    if service_cfg:
        session2 = create_odoo_session(service_cfg["url"], service_cfg["db"], service_cfg["username"], service_cfg["password"])
        resp2 = session2.post(f"{service_cfg['url']}/web/dataset/call_kw", json={
            "jsonrpc": "2.0",
            "method": "call",
            "params": {
                "model": "service.question.feedback",
                "method": "search_read",
                "args": [],
                "kwargs": {
                    "domain": [
                        ["processing_staff", "in", [service_cfg["user_id"]]],
                        ["state", "in", ["draft", "open"]],
                    ],
                    "fields": ["respondent", "system"],
                    "limit": 200,
                },
            },
        }).json()
        if "result" in resp2:
            for task in resp2["result"]:
                project_name = task["respondent"][1] if task.get("respondent") else (
                    task["system"][1] if task.get("system") else "未知專案"
                )
                stats[project_name] += 1
        else:
            print(f"[WARN] service.question.feedback 查詢失敗: {resp2.get('error')}")

    return stats


def _fill_table_header(table, row_idx, headers, color):
    for col_idx, text in enumerate(headers):
        cell = table.cell(row_idx, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = color


def _apply_table_style(table):
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Microsoft JhengHei"
                paragraph.font.size = Pt(11)
                paragraph.alignment = PP_ALIGN.CENTER


def _style_section_title(cell, text):
    """設定合併標題列的樣式（深藍底、白字、粗體）。"""
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(31, 73, 125)
    for para in cell.text_frame.paragraphs:
        para.font.name = "Microsoft JhengHei"
        para.font.bold = True
        para.font.size = Pt(13)
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER



def generate_ppt(start_date, end_date, project_stats, total_all_hours, pending_stats, output_path, memo_lines=None):
    dt_start = datetime.strptime(start_date, "%Y-%m-%d").strftime("%m/%d")
    dt_end = datetime.strptime(end_date, "%Y-%m-%d").strftime("%m/%d")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── 投影片標題 ──
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = f"過去兩周工作進度 {dt_start}~{dt_end}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Microsoft JhengHei"

    # ── 資料準備 ──
    left_data = [
        [proj, f"{data['hours']:.1f}hr ({data['hours'] / total_all_hours * 100:.1f}%)", str(len(data["tasks"]))]
        for proj, data in sorted(project_stats.items(), key=lambda x: x[1]["hours"], reverse=True)
    ]
    right_data = [
        [proj, str(count)]
        for proj, count in sorted(pending_stats.items(), key=lambda x: x[1], reverse=True)
    ]

    # ── 版面常數 ──
    ROW_H = Inches(0.35)
    TABLE_TOP = Inches(1.1)
    LEFT_X = Inches(0.4)
    RIGHT_X = Inches(5.1)
    TABLE_W_L = Inches(4.5)
    TABLE_W_R = Inches(4.5)
    HEADER_COLOR = RGBColor(79, 129, 189)

    # 每張表：section title(1) + column header(1) + data rows
    left_total_rows = len(left_data) + 2
    right_total_rows = len(right_data) + 2
    left_h = ROW_H * left_total_rows
    right_h = ROW_H * right_total_rows
    max_h = max(left_h, right_h)

    NOTES_TOP = TABLE_TOP + max_h + Inches(0.47)
    _memo = memo_lines or []
    NOTES_H = Inches(0.35 + 0.28 * max(1, len(_memo)))

    # ── 左表：過去兩周工作內容 ──
    lt_shape = slide.shapes.add_table(left_total_rows, 3, LEFT_X, TABLE_TOP, TABLE_W_L, left_h)
    lt = lt_shape.table
    lt.columns[0].width = Inches(2.2)
    lt.columns[1].width = Inches(1.5)
    lt.columns[2].width = Inches(0.8)
    lt.cell(0, 0).merge(lt.cell(0, 2))
    _style_section_title(lt.cell(0, 0), f"過去兩周工作內容 {dt_start}~{dt_end}")
    _fill_table_header(lt, 1, ["項目名稱", "總工時(hr)", "數量"], HEADER_COLOR)
    for i, row_data in enumerate(left_data):
        for j, text in enumerate(row_data):
            lt.cell(i + 2, j).text = text
    _apply_table_style(lt)
    _style_section_title(lt.cell(0, 0), f"過去兩周工作內容 {dt_start}~{dt_end}")

    # ── 右表：未來工作計畫 ──
    rt_shape = slide.shapes.add_table(right_total_rows, 2, RIGHT_X, TABLE_TOP, TABLE_W_R, right_h)
    rt = rt_shape.table
    rt.columns[0].width = Inches(3.3)
    rt.columns[1].width = Inches(1.2)
    rt.cell(0, 0).merge(rt.cell(0, 1))
    _style_section_title(rt.cell(0, 0), "未來工作計畫")
    _fill_table_header(rt, 1, ["項目名稱", "數量"], HEADER_COLOR)
    for i, row_data in enumerate(right_data):
        for j, text in enumerate(row_data):
            rt.cell(i + 2, j).text = text
    _apply_table_style(rt)
    _style_section_title(rt.cell(0, 0), "未來工作計畫")

    # ── 備註區塊（橫跨左右兩表）──
    notes_box = slide.shapes.add_textbox(LEFT_X, NOTES_TOP, RIGHT_X + TABLE_W_R - LEFT_X, NOTES_H)
    tf = notes_box.text_frame
    tf.word_wrap = True
    p_title = tf.paragraphs[0]
    p_title.text = "備註"
    p_title.font.name = "Microsoft JhengHei"
    p_title.font.bold = True
    p_title.font.size = Pt(12)
    for line in _memo:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(11)

    ppt_file = output_path / "工作週報.pptx"
    prs.save(ppt_file)
    print(f"[SUCCESS] PPT 簡報已成功動態生成至: {ppt_file}")


def main():
    script_dir = Path(__file__).parent
    output_path = script_dir

    odoo_cfg = load_config("odoo")
    service_cfg = load_config("odoo_service", optional=True)

    memo_path = script_dir / "memo.txt"
    memo_lines = memo_path.read_text(encoding="utf-8").splitlines() if memo_path.exists() else []

    start_date, end_date = get_previous_two_weeks_range()
    print(f"[INFO] 觸發日前兩週區間: {start_date} ~ {end_date}")

    # 左表資料：工時統計
    session = create_odoo_session(odoo_cfg["url"], odoo_cfg["db"], odoo_cfg["username"], odoo_cfg["password"])
    ts_resp = session.post(f"{odoo_cfg['url']}/web/dataset/call_kw", json={
        "jsonrpc": "2.0",
        "method": "call",
        "params": {
            "model": "account.analytic.line",
            "method": "search_read",
            "args": [],
            "kwargs": {
                "domain": [
                    ["user_id", "=", odoo_cfg["user_id"]],
                    ["date", ">=", start_date],
                    ["date", "<=", end_date],
                ],
                "fields": ["name", "date", "unit_amount", "project_id", "task_id"],
                "order": "date desc",
                "limit": 150,
            },
        },
    }).json()

    if "error" in ts_resp:
        print(f"[ERROR] 工時表查詢失敗: {ts_resp['error']}")
        return

    timesheets = ts_resp.get("result", [])
    if not timesheets:
        print(f"[INFO] 區間 {start_date} ~ {end_date} 內沒有任何工時紀錄。")
        return

    project_stats = defaultdict(lambda: {"hours": 0.0, "tasks": set()})
    total_all_hours = 0.0
    for ts in timesheets:
        if datetime.strptime(ts.get("date", ""), "%Y-%m-%d").date().weekday() in (5, 6):
            continue
        ts_hours = ts.get("unit_amount", 0.0)
        description = remove_images_only(ts.get("name")).strip()
        project_name = ts["project_id"][1] if ts.get("project_id") else "未知專案"
        task_name = ts["task_id"][1] if ts.get("task_id") else ""
        key = _resolve_pending_key(project_name, task_name)
        if key == "教育訓練" and "AI" in description:
            key = "AI流程調整"
        project_stats[key]["hours"] += ts_hours
        if description:
            project_stats[key]["tasks"].add(description)
        total_all_hours += ts_hours

    # 右表資料：未完成任務
    print("[INFO] 抓取未完成任務...")
    pending_stats = fetch_pending_tasks(odoo_cfg, service_cfg)

    generate_ppt(start_date, end_date, project_stats, total_all_hours, pending_stats, output_path, memo_lines)


if __name__ == "__main__":
    main()
